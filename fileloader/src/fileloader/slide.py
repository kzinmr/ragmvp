"""Slides parser.

Contains parsers for .pptx files.

Image captioning is not supported.
"""

from pathlib import Path
from typing import Any

from fsspec import AbstractFileSystem
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table

from fileloader.base import BaseReader
from fileloader.schema import Document


def extract_table_text(table: Table) -> str:
    """Extract text from a table."""
    table_text = "\nTable Contents:\n"
    for row in table.rows:
        row_text = [cell.text.strip() for cell in row.cells]
        table_text += " | ".join(row_text) + "\n"
    return table_text + "\n"


class PptxReader(BaseReader):
    """Powerpoint parser.

    Extract text, tables, and specify slides.
    """

    def load_data(
        self,
        file: Path,
        extra_info: dict[str, Any] | None = None,
        fs: AbstractFileSystem | None = None,
    ) -> list[Document]:
        """Parse .pptx file."""
        if extra_info is None:
            extra_info = {}

        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)

        results = []
        result = ""
        slides = [
            {"slide": slide, "page": i}
            for i, slide in enumerate(presentation.slides, 1)
        ]
        extra_info["total_pages"] = len(slides)
        extra_info["file_path"] = str(file)

        for d in slides:
            slide = d["slide"]
            extra_info["page_number"] = int(d["page"])

            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    pass

                if isinstance(shape, GraphicFrame) and shape.has_table:
                    result += extract_table_text(shape.table)

                if hasattr(shape, "text"):
                    result += f"{shape.text}\n"

            results.append(Document(text=result, metadata=extra_info))
            result = ""

        return results
